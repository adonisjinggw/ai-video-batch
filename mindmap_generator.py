#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RollRoll AI 项目功能思维导图生成器
使用 python-pptx 库生成专业的思维导图PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 配色方案 - Ocean Gradient
COLORS = {
    'primary': RGBColor(0x06, 0x5A, 0x82),      # 深蓝
    'secondary': RGBColor(0x1C, 0x72, 0x93),    # 青蓝
    'accent': RGBColor(0x02, 0xC3, 0x9A),       # 薄荷绿
    'dark': RGBColor(0x21, 0x29, 0x5C),         # 深夜蓝
    'light': RGBColor(0xE8, 0xF4, 0xF8),        # 浅蓝白
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'text': RGBColor(0x1E, 0x29, 0x3B),
    'text_light': RGBColor(0x64, 0x74, 0x8B),
}

def add_title_bar(slide, text, bg_color=None):
    """添加标题栏"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or COLORS['primary']
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    return shape

def add_text_box(slide, text, left, top, width, height, font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text']
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, title, content_lines, title_color=None):
    """添加卡片"""
    # 卡片背景
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['white']
    shape.line.fill.background()
    
    # 标题
    add_text_box(slide, title, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4),
                 font_size=14, bold=True, color=title_color or COLORS['primary'])
    
    # 内容
    content_text = '\n'.join(['• ' + line for line in content_lines])
    add_text_box(slide, content_text, left + Inches(0.1), top + Inches(0.5), width - Inches(0.2), height - Inches(0.6),
                 font_size=10, color=COLORS['text'])
    
    return shape

# ==================== Slide 1: 封面 ====================
slide_layout = prs.slide_layouts[6]  # 空白布局
slide1 = prs.slides.add_slide(slide_layout)

# 背景
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg.fill.solid()
bg.fill.fore_color.rgb = COLORS['dark']
bg.line.fill.background()

# 主标题
add_text_box(slide1, "🎬 RollRoll AI", Inches(0.5), Inches(1.5), Inches(9), Inches(1),
             font_size=54, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide1, "AI 多媒体创作平台功能全景", Inches(0.5), Inches(2.6), Inches(9), Inches(0.6),
             font_size=28, color=COLORS['accent'], align=PP_ALIGN.CENTER)

# 技术栈
add_text_box(slide1, "Serverless + Vanilla JS + Supabase", Inches(0.5), Inches(3.3), Inches(9), Inches(0.5),
             font_size=18, color=COLORS['text_light'], align=PP_ALIGN.CENTER)

# 域名
add_text_box(slide1, "www.rollroll.art | lossloop.cn", Inches(0.5), Inches(4.5), Inches(9), Inches(0.4),
             font_size=14, color=COLORS['text_light'], align=PP_ALIGN.CENTER)

# ==================== Slide 2: 项目总览 ====================
slide2 = prs.slides.add_slide(slide_layout)

# 背景
bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg2.fill.solid()
bg2.fill.fore_color.rgb = COLORS['light']
bg2.line.fill.background()

# 标题栏
add_title_bar(slide2, "🎯 项目总览")

# 核心定位卡片
card_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.0), Inches(9), Inches(0.9))
card_bg.fill.solid()
card_bg.fill.fore_color.rgb = COLORS['white']
card_bg.line.fill.background()

add_text_box(slide2, "一站式 AI 创意平台：视频生成 • 图像创作 • 音频制作 • 文本写作 • 多智能体协作",
             Inches(0.6), Inches(1.25), Inches(8.8), Inches(0.5),
             font_size=18, color=COLORS['text'], align=PP_ALIGN.CENTER)

# 三大核心特性
features = [
    ("⚡", "Serverless 架构", "Vercel 无服务器函数\n按需扩展，零运维"),
    ("🔐", "统一认证计费", "Supabase 认证 + 存储\n两阶段扣费机制"),
    ("🤖", "多模型集成", "10+ 视频模型\n10+ 图像模型")
]

for i, (icon, title, desc) in enumerate(features):
    x = Inches(0.5 + i * 3.1)
    
    # 卡片背景
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.9), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS['white']
    card.line.fill.background()
    
    add_text_box(slide2, icon, x, Inches(2.3), Inches(2.9), Inches(0.5),
                 font_size=28, color=COLORS['text'], align=PP_ALIGN.CENTER)
    add_text_box(slide2, title, x + Inches(0.1), Inches(2.8), Inches(2.7), Inches(0.35),
                 font_size=13, bold=True, color=COLORS['primary'], align=PP_ALIGN.CENTER)
    add_text_box(slide2, desc, x + Inches(0.1), Inches(3.2), Inches(2.7), Inches(0.7),
                 font_size=10, color=COLORS['text_light'], align=PP_ALIGN.CENTER)

# 技术栈说明
add_text_box(slide2, "技术栈：前端 (原生JS/CSS) → 后端 (Vercel Functions) → 数据库 (Supabase) → AI模型 (Sora2/Wan/Gemini/MJ等)",
             Inches(0.5), Inches(4.3), Inches(9), Inches(0.5),
             font_size=11, color=COLORS['text_light'], align=PP_ALIGN.CENTER)

# ==================== Slide 3: 页面结构 ====================
slide3 = prs.slides.add_slide(slide_layout)

bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg3.fill.solid()
bg3.fill.fore_color.rgb = COLORS['light']
bg3.line.fill.background()

add_title_bar(slide3, "📄 页面文件结构", bg_color=COLORS['secondary'])

# 主要入口页面
add_text_box(slide3, "主要入口页面", Inches(0.5), Inches(0.9), Inches(3), Inches(0.35),
             font_size=14, bold=True, color=COLORS['primary'])

main_pages = [
    ("index.html", "PC端主页面 - 视频批量生成工作台"),
    ("mobile.html", "移动端主页面 (734KB) - 功能最完整"),
    ("chat.html", "AI聊天界面 - 集成多智能体团队系统")
]

for i, (file, desc) in enumerate(main_pages):
    y = Inches(1.3 + i * 0.5)
    # 左侧色条
    bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.06), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()
    
    add_text_box(slide3, file, Inches(0.65), y, Inches(2), Inches(0.4),
                 font_size=11, bold=True, color=COLORS['text'])
    add_text_box(slide3, desc, Inches(2.7), y, Inches(2.3), Inches(0.4),
                 font_size=10, color=COLORS['text_light'])

# 功能子页面
add_text_box(slide3, "功能子页面", Inches(5.3), Inches(0.9), Inches(3), Inches(0.35),
             font_size=14, bold=True, color=COLORS['primary'])

sub_pages = [
    ("banana.html", "AI画图工具"),
    ("video-tools.html", "视频工具箱"),
    ("voice.html", "AI配音 (TTS)"),
    ("music.html", "AI音乐生成"),
    ("writing.html", "AI写作工具"),
    ("knolling.html", "Knolling拆解图")
]

for i, (file, desc) in enumerate(sub_pages):
    y = Inches(1.3 + i * 0.45)
    bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), y, Inches(0.05), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['secondary']
    bar.line.fill.background()
    
    add_text_box(slide3, file, Inches(5.45), y, Inches(2), Inches(0.35),
                 font_size=10, bold=True, color=COLORS['text'])
    add_text_box(slide3, desc, Inches(7.5), y, Inches(2), Inches(0.35),
                 font_size=10, color=COLORS['text_light'])

# 用户系统页面
add_text_box(slide3, "用户系统页面", Inches(0.5), Inches(3.0), Inches(3), Inches(0.35),
             font_size=14, bold=True, color=COLORS['primary'])

user_pages = [
    ("auth.html", "登录/注册"),
    ("user.html", "用户中心"),
    ("buy.html", "充值购买"),
    ("welcome-*.html", "欢迎页系列")
]

for i, (file, desc) in enumerate(user_pages):
    y = Inches(3.4 + i * 0.4)
    add_text_box(slide3, file, Inches(0.65), y, Inches(2), Inches(0.35),
                 font_size=10, bold=True, color=COLORS['text'])
    add_text_box(slide3, desc, Inches(2.7), y, Inches(2), Inches(0.35),
                 font_size=10, color=COLORS['text_light'])

# ==================== Slide 4: 核心JS模块 ====================
slide4 = prs.slides.add_slide(slide_layout)

bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg4.fill.solid()
bg4.fill.fore_color.rgb = COLORS['light']
bg4.line.fill.background()

add_title_bar(slide4, "🧩 核心JS模块架构", bg_color=COLORS['accent'])

# 技能系统卡片
add_card(slide4, Inches(0.3), Inches(0.9), Inches(3.0), Inches(2.0),
         "🎯 技能系统", [
             "skill-system.js",
             "SkillManager 管理器",
             "技能注册/执行/取消",
             "最大并发3控制",
             "skill-presets.js",
             "100+ 预置技能模板"
         ])

# 多智能体团队卡片
add_card(slide4, Inches(3.5), Inches(0.9), Inches(3.0), Inches(2.0),
         "🤖 多智能体团队", [
             "agent-team.js",
             "ToolRegistry 工具注册",
             "AgentTeam 协作调度",
             "agent-roles.js",
             "21个专业角色",
             "10个预设团队模板"
         ])

# API核心卡片
add_card(slide4, Inches(6.7), Inches(0.9), Inches(3.0), Inches(2.0),
         "🔌 API核心", [
             "api-core.js",
             "统一API调用封装",
             "视频/图片生成API",
             "TTS/OCR API",
             "resilient-api.js",
             "负载均衡/熔断降级"
         ])

# 基础设施卡片
add_card(slide4, Inches(0.3), Inches(3.1), Inches(4.6), Inches(1.4),
         "🏗️ 基础设施", [
             "billing.js - 两阶段扣费",
             "task-orchestrator.js - 任务调度器",
             "supabase-config.js - 认证配置"
         ])

# 其他模块卡片
add_card(slide4, Inches(5.1), Inches(3.1), Inches(4.6), Inches(1.4),
         "📚 其他模块", [
             "novel-engine.js - 长篇小说引擎",
             "prompt-templates.js - 100+ 提示词模板",
             "batch.js - PC端批量工作流"
         ])

# ==================== Slide 5: 技能系统详情 ====================
slide5 = prs.slides.add_slide(slide_layout)

bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg5.fill.solid()
bg5.fill.fore_color.rgb = COLORS['light']
bg5.line.fill.background()

add_title_bar(slide5, "🎯 技能系统详解")

# SkillManager卡片
add_card(slide5, Inches(0.3), Inches(0.9), Inches(4.5), Inches(1.7),
         "SkillManager 核心功能", [
             "register() - 技能注册",
             "execute() - 执行技能 (进度回调)",
             "cancel() - 取消执行",
             "最大并发控制 (默认3)",
             "一次性预扣费机制",
             "历史记录 & 收藏功能"
         ])

# SkillUI卡片
add_card(slide5, Inches(5.0), Inches(0.9), Inches(4.5), Inches(1.7),
         "SkillUI 界面组件", [
             "技能卡片列表渲染",
             "分类过滤 & 搜索",
             "参数表单自动生成",
             "支持类型: text/number/textarea/",
             "  select/checkbox/file/image",
             "进度面板显示"
         ])

# 技能分类标题
add_text_box(slide5, "技能分类 (7大类)", Inches(0.5), Inches(2.8), Inches(3), Inches(0.35),
             font_size=14, bold=True, color=COLORS['primary'])

# 技能分类标签
skill_categories = [
    ("🎬", "视频创作"),
    ("🎨", "图像创作"),
    ("✍️", "文本内容"),
    ("🎵", "音频制作"),
    ("🖼️", "设计工具"),
    ("🔧", "实用工具"),
    ("⚡", "自动化")
]

for i, (icon, name) in enumerate(skill_categories):
    x = Inches(0.4 + i * 1.35)
    tag = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(1.25), Inches(0.7))
    tag.fill.solid()
    tag.fill.fore_color.rgb = COLORS['secondary']
    tag.line.fill.background()
    
    add_text_box(slide5, icon, x, Inches(3.22), Inches(1.25), Inches(0.35),
                 font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide5, name, x, Inches(3.52), Inches(1.25), Inches(0.3),
                 font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER)

# 支持模型说明
add_text_box(slide5, "支持的AI模型", Inches(0.5), Inches(4.1), Inches(3), Inches(0.3),
             font_size=12, bold=True, color=COLORS['text'])
add_text_box(slide5, "视频: Wan2.6 / Grok Video 3 / Veo 3.1 / Vidu / Kling / Hailuo / LTX-Video",
             Inches(0.5), Inches(4.4), Inches(9), Inches(0.3),
             font_size=10, color=COLORS['text_light'])
add_text_box(slide5, "图像: Gemini Flash / Banana / 星梦画师 / MJ / ModelScope",
             Inches(0.5), Inches(4.7), Inches(9), Inches(0.3),
             font_size=10, color=COLORS['text_light'])

# ==================== Slide 6: 多智能体团队 ====================
slide6 = prs.slides.add_slide(slide_layout)

bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg6.fill.solid()
bg6.fill.fore_color.rgb = COLORS['light']
bg6.line.fill.background()

add_title_bar(slide6, "🤖 多智能体团队系统", bg_color=COLORS['secondary'])

# 核心组件标题
add_text_box(slide6, "核心组件架构", Inches(0.5), Inches(0.9), Inches(3), Inches(0.35),
             font_size=14, bold=True, color=COLORS['primary'])

# 三个核心组件
components = [
    ("ToolRegistry", "工具注册表", "映射 tool_id → API函数\n20+ 工具能力"),
    ("AgentTeam", "团队协作", "LLM驱动智能调度\n并行/串行编排"),
    ("AgentUI", "交互面板", "团队选择/配置/执行\n结果展示")
]

for i, (name, subtitle, desc) in enumerate(components):
    x = Inches(0.3 + i * 3.2)
    
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.0), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS['white']
    card.line.fill.background()
    
    add_text_box(slide6, name, x, Inches(1.4), Inches(3.0), Inches(0.35),
                 font_size=12, bold=True, color=COLORS['primary'], align=PP_ALIGN.CENTER)
    add_text_box(slide6, subtitle, x, Inches(1.7), Inches(3.0), Inches(0.3),
                 font_size=10, color=COLORS['text'], align=PP_ALIGN.CENTER)
    add_text_box(slide6, desc, x + Inches(0.1), Inches(2.0), Inches(2.8), Inches(0.45),
                 font_size=9, color=COLORS['text_light'], align=PP_ALIGN.CENTER)

# 21个角色
add_text_box(slide6, "21个专业智能体角色", Inches(0.5), Inches(2.7), Inches(4), Inches(0.35),
             font_size=14, bold=True, color=COLORS['primary'])

roles = [
    "👔 项目总监", "✍️ 文案策划", "🎨 视觉设计师", "🎬 视频制作",
    "💡 品牌顾问", "🧸 角色设计师", "🎤 配音师", "🎥 导演",
    "🖼️ 分镜大师", "✏️ 漫画家", "🎵 音乐制作人", "🧊 3D建模师"
]

for i, role in enumerate(roles):
    x = Inches(0.4 + (i % 4) * 2.4)
    y = Inches(3.1 + (i // 4) * 0.45)
    add_text_box(slide6, role, x, y, Inches(2.3), Inches(0.4),
                 font_size=11, color=COLORS['text'])

# 10个团队模板
add_text_box(slide6, "10个预设团队模板", Inches(0.5), Inches(4.5), Inches(4), Inches(0.3),
             font_size=12, bold=True, color=COLORS['text'])
add_text_box(slide6, "品牌全案 | 短视频 | IP设计 | 电商 | 音频制作 | 有声小说 | 漫画创作 | 小说转短剧 | 影视制作 | 混元3D建模",
             Inches(0.5), Inches(4.8), Inches(9), Inches(0.3),
             font_size=10, color=COLORS['text_light'])

# ==================== Slide 7: API端点 ====================
slide7 = prs.slides.add_slide(slide_layout)

bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg7.fill.solid()
bg7.fill.fore_color.rgb = COLORS['light']
bg7.line.fill.background()

add_title_bar(slide7, "🔌 Serverless API 端点", bg_color=COLORS['accent'])

# 表头
header_bg = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.9), Inches(9.4), Inches(0.4))
header_bg.fill.solid()
header_bg.fill.fore_color.rgb = COLORS['primary']
header_bg.line.fill.background()

add_text_box(slide7, "API文件", Inches(0.4), Inches(0.95), Inches(2.5), Inches(0.35),
             font_size=11, bold=True, color=COLORS['white'])
add_text_box(slide7, "功能描述", Inches(3.0), Inches(0.95), Inches(4.5), Inches(0.35),
             font_size=11, bold=True, color=COLORS['white'])
add_text_box(slide7, "计费", Inches(7.6), Inches(0.95), Inches(2), Inches(0.35),
             font_size=11, bold=True, color=COLORS['white'])

# API列表
apis = [
    ("sora2.js", "Sora2 视频生成代理", "7~14胶片"),
    ("banana2.js", "Banana2 图片生成代理", "4~10胶片"),
    ("yunwu.js", "云雾AI多模态代理", "按配置"),
    ("suno.js", "Suno 音乐生成代理", "8胶片/首"),
    ("modelscope.js", "ModelScope 免费API", "0~3胶片"),
    ("supabase-proxy.js", "数据库/认证代理", "-"),
    ("proxy.js", "通用代理 + IP限流", "-"),
    ("writer-llm.js", "写作 LLM", "-")
]

for i, (file, desc, cost) in enumerate(apis):
    y = Inches(1.35 + i * 0.42)
    
    # 行背景
    row_bg = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), y, Inches(9.4), Inches(0.4))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = COLORS['white'] if i % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFC)
    row_bg.line.fill.background()
    
    add_text_box(slide7, file, Inches(0.4), y + Inches(0.05), Inches(2.5), Inches(0.3),
                 font_size=10, bold=True, color=COLORS['secondary'])
    add_text_box(slide7, desc, Inches(3.0), y + Inches(0.05), Inches(4.5), Inches(0.3),
                 font_size=10, color=COLORS['text'])
    add_text_box(slide7, cost, Inches(7.6), y + Inches(0.05), Inches(2), Inches(0.3),
                 font_size=10, color=COLORS['text_light'])

# 计费说明
add_text_box(slide7, "胶片计费单位: 1胶片 = 10 units | 通过 /api/supabase-proxy 的 consume/recharge action 实现",
             Inches(0.3), Inches(4.8), Inches(9.4), Inches(0.3),
             font_size=10, color=COLORS['text_light'])

# ==================== Slide 8: 计费系统 ====================
slide8 = prs.slides.add_slide(slide_layout)

bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg8.fill.solid()
bg8.fill.fore_color.rgb = COLORS['light']
bg8.line.fill.background()

add_title_bar(slide8, "💰 统一计费系统")

# 两阶段扣费流程标题
add_text_box(slide8, "两阶段扣费流程", Inches(0.5), Inches(0.9), Inches(3), Inches(0.35),
             font_size=14, bold=True, color=COLORS['primary'])

# 流程步骤
steps = [
    ("1", "reserveFilm()", "预扣费冻结", COLORS['secondary']),
    ("2", "API调用", "执行请求", COLORS['accent']),
    ("3a", "commitFilm()", "确认扣费", RGBColor(0x22, 0xC5, 0x5E)),
    ("3b", "releaseFilm()", "释放冻结", RGBColor(0xEF, 0x44, 0x44))
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.4)
    
    # 圆形
    circle = slide8.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), Inches(1.4), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    add_text_box(slide8, num, x + Inches(0.5), Inches(1.55), Inches(0.7), Inches(0.4),
                 font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide8, title, x, Inches(2.2), Inches(2.0), Inches(0.35),
                 font_size=11, bold=True, color=COLORS['text'], align=PP_ALIGN.CENTER)
    add_text_box(slide8, desc, x, Inches(2.5), Inches(2.0), Inches(0.3),
                 font_size=10, color=COLORS['text_light'], align=PP_ALIGN.CENTER)

# 视频模型计费表
add_text_box(slide8, "视频模型计费表", Inches(0.5), Inches(3.0), Inches(3), Inches(0.35),
             font_size=13, bold=True, color=COLORS['primary'])

video_costs = [
    ("模型", "胶片"),
    ("Wan2.6 720p 5s", "3"),
    ("Wan2.6 720p 10s 有声", "7"),
    ("Wan2.6 1080p 15s 有声", "21"),
    ("Grok Video 3 6s/15s", "5~12"),
    ("Veo 3.1", "30"),
    ("Kling 2.5 720p 5s", "6")
]

for i, (model, cost) in enumerate(video_costs):
    y = Inches(3.4 + i * 0.35)
    bg_color = COLORS['primary'] if i == 0 else COLORS['white']
    text_color = COLORS['white'] if i == 0 else COLORS['text']
    
    row = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), y, Inches(4.5), Inches(0.32))
    row.fill.solid()
    row.fill.fore_color.rgb = bg_color
    row.line.fill.background()
    
    add_text_box(slide8, model, Inches(0.4), y + Inches(0.02), Inches(3.0), Inches(0.28),
                 font_size=9, bold=(i==0), color=text_color)
    add_text_box(slide8, cost, Inches(3.5), y + Inches(0.02), Inches(1.2), Inches(0.28),
                 font_size=9, color=text_color)

# TTS计费表
add_text_box(slide8, "TTS配音计费", Inches(5.3), Inches(3.0), Inches(3), Inches(0.35),
             font_size=13, bold=True, color=COLORS['primary'])

tts_costs = [
    ("引擎", "胶片"),
    ("DubbingX", "2"),
    ("Kling TTS", "2"),
    ("Gemini Flash", "1"),
    ("Gemini Pro", "3")
]

for i, (engine, cost) in enumerate(tts_costs):
    y = Inches(3.4 + i * 0.35)
    bg_color = COLORS['primary'] if i == 0 else COLORS['white']
    text_color = COLORS['white'] if i == 0 else COLORS['text']
    
    row = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), y, Inches(4.3), Inches(0.32))
    row.fill.solid()
    row.fill.fore_color.rgb = bg_color
    row.line.fill.background()
    
    add_text_box(slide8, engine, Inches(5.3), y + Inches(0.02), Inches(2.8), Inches(0.28),
                 font_size=9, bold=(i==0), color=text_color)
    add_text_box(slide8, cost, Inches(8.2), y + Inches(0.02), Inches(1.2), Inches(0.28),
                 font_size=9, color=text_color)

# ==================== Slide 9: 基础设施特性 ====================
slide9 = prs.slides.add_slide(slide_layout)

bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg9.fill.solid()
bg9.fill.fore_color.rgb = COLORS['light']
bg9.line.fill.background()

add_title_bar(slide9, "🏗️ 基础设施特性", bg_color=COLORS['dark'])

# 弹性API网关
add_card(slide9, Inches(0.3), Inches(0.9), Inches(4.5), Inches(1.8),
         "弹性API网关 (resilient-api.js)", [
             "负载均衡 - 多节点自动分发",
             "熔断机制 - 连续3次失败触发熔断",
             "健康检查 - 定期探测节点状态",
             "成本优化 - 根据价格选择最优节点",
             "智能降级 - 付费→免费自动切换"
         ])

# 任务调度器
add_card(slide9, Inches(5.0), Inches(0.9), Inches(4.5), Inches(1.8),
         "任务调度器 (task-orchestrator.js)", [
             "优先级队列 - URGENT/HIGH/NORMAL/LOW",
             "智能重试 - 指数退避，最大3次",
             "断点续传 - localStorage持久化",
             "预估时间 - 基于历史数据动态预估",
             "并发控制 - 最大并发3"
         ])

# 认证系统
add_card(slide9, Inches(0.3), Inches(2.9), Inches(4.5), Inches(1.3),
         "认证系统 (supabase-config.js)", [
             "Supabase Auth 集成",
             "本地缓存快速恢复 (30分钟有效)",
             "多重存储备份"
         ])

# 移动端适配
add_card(slide9, Inches(5.0), Inches(2.9), Inches(4.5), Inches(1.3),
         "移动端适配", [
             "自动跳转移动端 → mobile.html",
             "响应式设计",
             "触控事件支持 (click + touchend)"
         ])

# ==================== Slide 10: 功能总结 ====================
slide10 = prs.slides.add_slide(slide_layout)

bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg10.fill.solid()
bg10.fill.fore_color.rgb = COLORS['dark']
bg10.line.fill.background()

# 标题
add_text_box(slide10, "🎬 功能全景总结", Inches(0.5), Inches(0.4), Inches(9), Inches(0.6),
             font_size=32, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

# 六大核心功能
core_features = [
    ("🎬", "视频创作", "文生视频/图生视频\n10+ 模型支持\n批量生成队列"),
    ("🎨", "图像创作", "多模型图片生成\n风格迁移/图生图\nIP角色设计"),
    ("🎙️", "音频制作", "TTS配音(多引擎)\n多角色配音\nAI音乐生成"),
    ("✍️", "文本创作", "剧本/文案生成\n长篇小说引擎\n剧情一致性校验"),
    ("🤖", "智能体协作", "21个专业角色\n10个团队模板\n自动任务调度"),
    ("🧩", "技能系统", "100+ 预置技能\n自定义技能\n进度跟踪")
]

for i, (icon, title, items) in enumerate(core_features):
    x = Inches(0.4 + (i % 3) * 3.15)
    y = Inches(1.1 + (i // 3) * 1.9)
    
    # 卡片背景
    card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.0), Inches(1.7))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x30, 0x38, 0x68)  # 深蓝半透明
    card.line.color.rgb = COLORS['accent']
    card.line.width = Pt(1)
    
    add_text_box(slide10, icon, x, y + Inches(0.1), Inches(3.0), Inches(0.4),
                 font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide10, title, x + Inches(0.1), y + Inches(0.5), Inches(2.8), Inches(0.35),
                 font_size=13, bold=True, color=COLORS['accent'], align=PP_ALIGN.CENTER)
    add_text_box(slide10, items, x + Inches(0.1), y + Inches(0.85), Inches(2.8), Inches(0.75),
                 font_size=9, color=COLORS['white'], align=PP_ALIGN.CENTER)

# 技术栈
add_text_box(slide10, "技术栈: 前端(原生JS) + 后端(Vercel Functions) + 数据库(Supabase) + AI(Sora2/Gemini/MJ等)",
             Inches(0.5), Inches(5.0), Inches(9), Inches(0.3),
             font_size=10, color=COLORS['text_light'], align=PP_ALIGN.CENTER)

# 保存文件
output_path = "j:/123pan/13998416173/NanoNoPort/ai-video-batch/RollRoll-AI-项目功能思维导图.pptx"
prs.save(output_path)
print(f"[OK] PPT generated: {output_path}")
